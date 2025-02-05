from pptx import Presentation
from typing import List, Dict, Any, Optional
import logging
from dataclasses import dataclass
from openai import OpenAI
import os

@dataclass
class TextFormat:
    font_name: Optional[str]
    font_size: Optional[int]
    bold: Optional[bool] 
    italic: Optional[bool]
    color: Optional[str]

@dataclass
class TextElement:
    slide_idx: int
    shape_idx: int
    original_text: str
    translated_text: Optional[str]
    format: TextFormat

class PPTTranslator:
    def __init__(self, api_key: str):
        """Initialize the PPT translator with DeepSeek API key"""
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://api.deepseek.com"  # DeepSeek API endpoint
        )
        
        # Configure logging
        self.logger = logging.getLogger(__name__)
        if not self.logger.handlers:
            handler = logging.StreamHandler()
            formatter = logging.Formatter(
                '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
            )
            handler.setFormatter(formatter)
            self.logger.addHandler(handler)
            self.logger.setLevel(logging.INFO)
    
    def extract_text_elements(self, ppt_path: str) -> List[TextElement]:
        """Extract text and formatting from PPT file"""
        prs = Presentation(ppt_path)
        elements = []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    text_frame = shape.text_frame
                    
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 修改颜色获取逻辑
                            color = None
                            try:
                                if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                    color = run.font.color.rgb
                            except Exception as e:
                                self.logger.warning(f"Failed to get color: {str(e)}")
                            
                            format = TextFormat(
                                font_name=run.font.name,
                                font_size=run.font.size,
                                bold=run.font.bold,
                                italic=run.font.italic,
                                color=color  # 使用处理后的颜色值
                            )
                            
                            element = TextElement(
                                slide_idx=slide_idx,
                                shape_idx=shape_idx,
                                original_text=run.text,
                                translated_text=None,
                                format=format
                            )
                            elements.append(element)
        
        return elements
    
    def translate_text(self, text: str, target_lang: str) -> str:
        """Translate text using DeepSeek API"""
        try:
            response = self.client.chat.completions.create(
                model="deepseek-chat",
                messages=[
                    {"role": "system", "content": f"You are a translator. Translate the following text to {target_lang}. Only return the translation, no explanations."},
                    {"role": "user", "content": text}
                ],
                stream=False
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            self.logger.error(f"Translation failed: {str(e)}")
            return text
    
    def translate_elements(self, elements: List[TextElement], target_lang: str) -> List[TextElement]:
        """Translate text elements to target language"""
        for element in elements:
            try:
                if element.original_text.strip():  # Only translate non-empty text
                    element.translated_text = self.translate_text(
                        element.original_text,
                        target_lang
                    )
                else:
                    element.translated_text = element.original_text
            except Exception as e:
                self.logger.error(f"Translation failed for text: {element.original_text}")
                self.logger.error(str(e))
                element.translated_text = element.original_text
                
        return elements
    
    def replace_text(self, ppt_path: str, elements: List[TextElement], output_path: str):
        """Replace original text with translations while preserving formatting"""
        prs = Presentation(ppt_path)
        
        # 修改文本匹配逻辑
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, "text_frame"):
                    continue
                    
                text_frame = shape.text_frame
                
                # 处理每个段落
                for paragraph in text_frame.paragraphs:
                    # 找到匹配的翻译
                    for element in elements:
                        if (element.slide_idx == slide_idx and 
                            element.shape_idx == shape_idx and 
                            element.original_text.strip() in paragraph.text):
                            
                            # 处理段落中的每个文本运行
                            for run in paragraph.runs:
                                if element.original_text.strip() in run.text:
                                    # 保存原始格式
                                    font = run.font
                                    original_format = {
                                        'name': font.name,
                                        'size': font.size,
                                        'bold': font.bold,
                                        'italic': font.italic,
                                        'color': font.color.rgb if font.color else None
                                    }
                                    
                                    # 替换文本
                                    run.text = run.text.replace(
                                        element.original_text.strip(),
                                        element.translated_text.strip()
                                    )
                                    
                                    # 重新应用格式
                                    try:
                                        if original_format['name']:
                                            run.font.name = original_format['name']
                                        if original_format['size']:
                                            run.font.size = original_format['size']
                                        if original_format['bold'] is not None:
                                            run.font.bold = original_format['bold']
                                        if original_format['italic'] is not None:
                                            run.font.italic = original_format['italic']
                                        if original_format['color']:
                                            if not run.font.color:
                                                run.font.color.rgb = original_format['color']
                                            elif hasattr(run.font.color, 'rgb'):
                                                run.font.color.rgb = original_format['color']
                                    except Exception as e:
                                        self.logger.warning(f"Failed to restore format: {str(e)}")
        
        # 保存修改后的演示文稿
        try:
            prs.save(output_path)
        except Exception as e:
            self.logger.error(f"Failed to save presentation: {str(e)}")
            raise
    
    def translate_ppt(self, input_path: str, output_path: str, target_lang: str = "zh-CN", progress_callback=None):
        """
        Translate PowerPoint presentation while preserving formatting.
        
        Args:
            input_path: Path to input PPT file
            output_path: Path to save translated PPT file
            target_lang: Target language code (default: zh-CN)
            progress_callback: Callback function for progress updates
            
        Raises:
            FileNotFoundError: If input file doesn't exist
            PermissionError: If can't write to output path
            Exception: For other errors during translation
        """
        try:
            # Validate input file
            if not os.path.exists(input_path):
                raise FileNotFoundError(f"Input file not found: {input_path}")
            
            self.logger.info(f"Starting translation of {input_path} to {target_lang}")
            
            # Step 1: Extract text and formatting
            self.logger.info("Extracting text and formatting...")
            elements = self.extract_text_elements(input_path)
            total_elements = len(elements)
            self.logger.info(f"Found {total_elements} text elements")
            
            # Step 2: Translate text
            self.logger.info("Translating text...")
            translated_elements = []
            
            for i, element in enumerate(elements, 1):
                try:
                    if element.original_text.strip():
                        element.translated_text = self.translate_text(
                            element.original_text,
                            target_lang
                        )
                    else:
                        element.translated_text = element.original_text
                    
                    translated_elements.append(element)
                    
                    # 调用进度回调
                    if progress_callback:
                        progress_callback(i, total_elements)
                    
                except Exception as e:
                    self.logger.error(f"Failed to translate text: {element.original_text}")
                    self.logger.error(str(e))
                    element.translated_text = element.original_text
                    translated_elements.append(element)
            
            # Step 3: Replace text while preserving formatting
            self.logger.info("Replacing text in PPT...")
            self.replace_text(input_path, translated_elements, output_path)
            
            self.logger.info(f"Translation completed. Output saved to: {output_path}")
            
        except Exception as e:
            self.logger.error(f"Translation failed: {str(e)}")
            raise 
